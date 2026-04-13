"""
WorkTrack AI — Visual Manager Presentation
9 slides, dark theme, icon-driven, minimal text, visual diagrams
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
import copy, math

# ── Palette ───────────────────────────────────────────────────────────────────
BG      = RGBColor(0x0A, 0x14, 0x28)   # deep navy
CARD    = RGBColor(0x12, 0x24, 0x3A)   # card dark
CARD2   = RGBColor(0x1A, 0x30, 0x4D)   # lighter card
BLUE    = RGBColor(0x00, 0xB4, 0xFF)   # electric blue
VIOLET  = RGBColor(0x7C, 0x3A, 0xED)   # violet
GREEN   = RGBColor(0x10, 0xB9, 0x81)   # emerald
AMBER   = RGBColor(0xF5, 0x9E, 0x0B)   # amber
RED     = RGBColor(0xEF, 0x44, 0x44)   # red accent
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY   = RGBColor(0xA0, 0xB8, 0xD0)
DGRAY   = RGBColor(0x1C, 0x38, 0x5A)
TEAL    = RGBColor(0x06, 0xB6, 0xD4)

W = Inches(13.33)
H = Inches(7.5)

# ── Core helpers ──────────────────────────────────────────────────────────────
def solid_fill(shape, rgb):
    fill = shape.fill; fill.solid(); fill.fore_color.rgb = rgb

def no_line(shape):
    shape.line.fill.background()

def rect(slide, l, t, w, h, rgb, line=True):
    s = slide.shapes.add_shape(1, l, t, w, h)
    solid_fill(s, rgb)
    if line: no_line(s)
    return s

def oval(slide, l, t, w, h, rgb):
    from pptx.util import Emu
    s = slide.shapes.add_shape(9, l, t, w, h)   # MSO_SHAPE.OVAL = 9
    solid_fill(s, rgb); no_line(s)
    return s

def txt(slide, text, l, t, w, h, size=14, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, italic=False, wrap=True):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = wrap
    p  = tf.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = text
    run.font.size = Pt(size); run.font.bold = bold
    run.font.italic = italic; run.font.color.rgb = color
    run.font.name = "Segoe UI"
    return tb

def bg(slide, prs, rgb=None):
    rgb = rgb or BG
    fill = slide.background.fill; fill.solid(); fill.fore_color.rgb = rgb

def icon(slide, symbol, l, t, size=28, color=WHITE, bold=True):
    w = Inches(0.8); h = Inches(0.7)
    txt(slide, symbol, l, t, w, h, size=size, bold=bold, color=color, align=PP_ALIGN.CENTER)

def hbar(slide, t, color=DGRAY, h_in=0.015):
    rect(slide, Inches(0.5), t, Inches(12.33), Inches(h_in), color)

def top_stripe(slide, color, h_in=0.055):
    rect(slide, 0, 0, W, Inches(h_in), color)

def bottom_bar(slide, color, text, text_color=WHITE, size=9.5):
    rect(slide, 0, H - Inches(0.42), W, Inches(0.42), color)
    txt(slide, text, Inches(0.5), H - Inches(0.4), Inches(12.3), Inches(0.38),
        size=size, color=text_color, align=PP_ALIGN.CENTER)

def badge(slide, label, l, t, color, w_in=1.4, h_in=0.32, size=10):
    rect(slide, l, t, Inches(w_in), Inches(h_in), color)
    txt(slide, label, l, t, Inches(w_in), Inches(h_in),
        size=size, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

def circle_num(slide, num, cx, cy, r, color):
    """Draw a circle badge with number at center (cx,cy)."""
    oval(slide, cx - r, cy - r, r*2, r*2, color)
    txt(slide, str(num), cx - r, cy - r*1.1, r*2, r*2,
        size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

def card(slide, l, t, w, h, color, accent_color=None, top_bar=True):
    r = rect(slide, l, t, w, h, color)
    if top_bar and accent_color:
        rect(slide, l, t, w, Inches(0.055), accent_color)
    return r

def arrow_right(slide, x, y, length=Inches(0.4)):
    """Horizontal right arrow as a text box."""
    txt(slide, "➜", x, y - Inches(0.25), length, Inches(0.5),
        size=18, bold=True, color=BLUE, align=PP_ALIGN.CENTER)

def arrow_down(slide, x, y):
    txt(slide, "▼", x, y, Inches(0.4), Inches(0.35),
        size=14, bold=True, color=BLUE, align=PP_ALIGN.CENTER)

def slide_title(slide, title, subtitle=None, title_color=WHITE):
    txt(slide, title, Inches(0.5), Inches(0.22), Inches(12), Inches(0.6),
        size=28, bold=True, color=title_color)
    if subtitle:
        txt(slide, subtitle, Inches(0.5), Inches(0.82), Inches(12), Inches(0.38),
            size=13, color=LGRAY, italic=True)
    hbar(slide, Inches(1.28))


# ═════════════════════════════════════════════════════════════════════════════
#  SLIDE 1 — COVER
# ═════════════════════════════════════════════════════════════════════════════
def s1_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, prs, BG)

    # Left dark panel
    rect(slide, 0, 0, Inches(7.2), H, CARD)
    # Accent diagonal stripe (overlapping)
    rect(slide, Inches(7.0), 0, Inches(0.22), H, BLUE)

    # Big decorative circle (background art)
    oval(slide, Inches(-1.5), Inches(-1.5), Inches(6), Inches(6), CARD2)
    oval(slide, Inches(0.2), Inches(4.8), Inches(3), Inches(3), CARD2)

    # AI tag
    badge(slide, "  AI-POWERED  ", Inches(0.5), Inches(0.9), VIOLET, w_in=1.8)

    # Main title
    txt(slide, "WorkTrack AI", Inches(0.5), Inches(1.5), Inches(6.3), Inches(1.3),
        size=52, bold=True, color=WHITE)

    # Tagline
    txt(slide, "Natural Language Work Tracking\nPowered by Claude Sonnet 4.6",
        Inches(0.5), Inches(2.9), Inches(6.3), Inches(0.9),
        size=17, color=LGRAY, italic=True)

    hbar(slide, Inches(4.0), color=DGRAY)

    # Three stat boxes
    stats = [("6", "Phases\nDelivered", BLUE), ("218", "Tests\nPassing", GREEN), ("3", "User\nRoles", VIOLET)]
    for i, (val, lbl, col) in enumerate(stats):
        lx = Inches(0.5 + i * 2.15)
        rect(slide, lx, Inches(4.2), Inches(1.9), Inches(1.1), CARD2)
        rect(slide, lx, Inches(4.2), Inches(1.9), Inches(0.04), col)
        txt(slide, val, lx, Inches(4.25), Inches(1.9), Inches(0.55),
            size=26, bold=True, color=col, align=PP_ALIGN.CENTER)
        txt(slide, lbl, lx, Inches(4.82), Inches(1.9), Inches(0.45),
            size=9.5, color=LGRAY, align=PP_ALIGN.CENTER)

    # Tech pills row
    techs = ["FastAPI", "LangChain", "LangGraph", "ChromaDB", "Streamlit"]
    cols_t = [BLUE, VIOLET, GREEN, AMBER, TEAL]
    for i, (tech, col) in enumerate(zip(techs, cols_t)):
        lx = Inches(0.5 + i * 1.34)
        rect(slide, lx, Inches(5.75), Inches(1.2), Inches(0.28), col)
        txt(slide, tech, lx, Inches(5.75), Inches(1.2), Inches(0.28),
            size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Right panel — problem teaser
    txt(slide, "THE CHALLENGE", Inches(7.55), Inches(1.1), Inches(5.4), Inches(0.38),
        size=11, bold=True, color=BLUE)

    problems = [
        ("⏱", "Form fatigue", "Rigid standup forms waste time daily"),
        ("🔍", "Lost context", "Work history in emails — not searchable"),
        ("📊", "Blind spots", "Managers lack real-time team visibility"),
        ("🔗", "Siloed data", "Tickets, meetings, learning — all separate"),
    ]
    for i, (ico, title, desc) in enumerate(problems):
        top = Inches(1.6 + i * 1.35)
        rect(slide, Inches(7.55), top, Inches(5.35), Inches(1.15), CARD2)
        rect(slide, Inches(7.55), top, Inches(0.05), Inches(1.15), AMBER)
        txt(slide, ico,  Inches(7.65), top + Inches(0.28), Inches(0.6), Inches(0.55),
            size=20, align=PP_ALIGN.CENTER)
        txt(slide, title, Inches(8.35), top + Inches(0.07), Inches(4.4), Inches(0.38),
            size=13, bold=True, color=WHITE)
        txt(slide, desc,  Inches(8.35), top + Inches(0.47), Inches(4.4), Inches(0.58),
            size=10.5, color=LGRAY)

    bottom_bar(slide, VIOLET, "Prototype · Python · LLM · Streamlit · SQLite + ChromaDB")


# ═════════════════════════════════════════════════════════════════════════════
#  SLIDE 2 — HOW IT WORKS (3-step visual pipeline)
# ═════════════════════════════════════════════════════════════════════════════
def s2_pipeline(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, prs)
    top_stripe(slide, BLUE)
    slide_title(slide, "How It Works", "Three-step AI pipeline — from plain text to structured insight")

    steps = [
        ("📝", "01", "Natural Language\nInput", BLUE,
         ["Employee types update in plain English",
          '"Fixed bug INC-231 (2h, done). Standup 30m."',
          "Date picker defaults to today",
          "No dropdowns, no forms, no structure needed",
          "⚡ Fallback: if LLM unavailable, user fills manually"]),
        ("🤖", "02", "AI Extraction\nPipeline", VIOLET,
         ["Claude Sonnet 4.6 · temperature = 0",
          "Splits compound tasks → separate items",
          "Resolves relative dates ('yesterday' → date)",
          "Normalises status (complete → done)",
          "Filters non-work noise automatically",
          "OutputFixingParser auto-retries bad JSON"]),
        ("✅", "03", "Human Review\n& Save", GREEN,
         ["Editable preview table with confidence scores",
          "Clarification prompts only for missing fields",
          "Warning if total hours > 12",
          "Confirm → SQLite commit first",
          "ChromaDB upsert (semantic index updated)",
          "← Edit button: revise without losing raw text"]),
    ]

    for i, (ico, num, title, col, bullets) in enumerate(steps):
        lx = Inches(0.32 + i * 4.35)
        ty = Inches(1.42)
        card(slide, lx, ty, Inches(4.15), Inches(5.65), CARD, col)

        # Big number + icon side by side
        rect(slide, lx + Inches(0.15), ty + Inches(0.18),
             Inches(0.65), Inches(0.65), col)
        txt(slide, num, lx + Inches(0.15), ty + Inches(0.18),
            Inches(0.65), Inches(0.65),
            size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        txt(slide, ico, lx + Inches(0.88), ty + Inches(0.15),
            Inches(0.7), Inches(0.7), size=26, align=PP_ALIGN.CENTER)
        txt(slide, title, lx + Inches(1.6), ty + Inches(0.17),
            Inches(2.4), Inches(0.7), size=13, bold=True, color=col)

        hbar(slide, ty + Inches(1.0), color=DGRAY)

        for j, b in enumerate(bullets):
            txt(slide, f"  {b}", lx + Inches(0.18), ty + Inches(1.1 + j * 0.74),
                Inches(3.82), Inches(0.68), size=10.5, color=LGRAY if j > 0 else WHITE)

        # Arrow between cards
        if i < 2:
            ax = lx + Inches(4.15) + Inches(0.02)
            ay = ty + Inches(2.7)
            txt(slide, "➜", ax, ay, Inches(0.31), Inches(0.4),
                size=18, bold=True, color=col, align=PP_ALIGN.CENTER)

    bottom_bar(slide, CARD,
               "SQLite = source of truth  ·  Re-submit soft-deletes old log + re-extracts  ·  "
               "Identity always from JWT — never from LLM")


# ═════════════════════════════════════════════════════════════════════════════
#  SLIDE 3 — AI EXTRACTION DEEP DIVE (visual before → after)
# ═════════════════════════════════════════════════════════════════════════════
def s3_extraction(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, prs)
    top_stripe(slide, VIOLET)
    slide_title(slide, "Smart Extraction Engine",
                "LangChain · Claude Sonnet 4.6 · Pydantic v2 · 8 built-in rules")

    # ── LEFT: Raw input ──────────────────────────────────────────────────────
    card(slide, Inches(0.32), Inches(1.42), Inches(3.6), Inches(5.65), CARD, AMBER)
    txt(slide, "📥  RAW INPUT", Inches(0.47), Inches(1.55), Inches(3.3), Inches(0.4),
        size=11, bold=True, color=AMBER)
    raw_msg = (
        '"Fixed critical login bug\n'
        ' INC-231 — took about 2h,\n'
        ' done now. Sprint planning\n'
        ' meeting this morning.\n'
        ' Reviewed Alice\'s PR for\n'
        ' data pipeline, still going.\n'
        ' Power cut killed half\n'
        ' the afternoon."'
    )
    txt(slide, raw_msg, Inches(0.42), Inches(2.0), Inches(3.3), Inches(4.8),
        size=11, color=WHITE, italic=True)

    # Arrow
    txt(slide, "➜", Inches(3.98), Inches(4.0), Inches(0.5), Inches(0.55),
        size=24, bold=True, color=VIOLET, align=PP_ALIGN.CENTER)

    # ── MIDDLE: Pipeline rules ────────────────────────────────────────────────
    card(slide, Inches(4.55), Inches(1.42), Inches(4.2), Inches(5.65), CARD, VIOLET)
    txt(slide, "⚙  EXTRACTION RULES", Inches(4.7), Inches(1.55), Inches(3.9), Inches(0.4),
        size=11, bold=True, color=VIOLET)

    rules = [
        ("✂️", "Split",      "Compound tasks → separate items"),
        ("📅", "Dates",      "Relative → absolute (yesterday)"),
        ("🔄", "Normalise",  "complete→done · waiting→blocked"),
        ("🚫", "Filter",     "Non-work noise skipped entirely"),
        ("❓", "Clarify",    "Only when hours AND status missing"),
        ("🏷", "Categories", "10 types: ticket, meeting, review…"),
        ("💯", "Confidence", "Score 0.0–1.0 per extracted item"),
        ("🔧", "Auto-fix",   "Haiku retries malformed JSON once"),
    ]
    for j, (ico, label, desc) in enumerate(rules):
        ty = Inches(2.05 + j * 0.63)
        rect(slide, Inches(4.62), ty, Inches(4.0), Inches(0.56), CARD2)
        txt(slide, ico,   Inches(4.7),  ty + Inches(0.1), Inches(0.4), Inches(0.38), size=14)
        txt(slide, label, Inches(5.15), ty + Inches(0.06), Inches(1.0), Inches(0.38),
            size=10, bold=True, color=WHITE)
        txt(slide, desc,  Inches(6.2),  ty + Inches(0.06), Inches(2.2), Inches(0.38),
            size=9.5, color=LGRAY)

    # Arrow
    txt(slide, "➜", Inches(8.82), Inches(4.0), Inches(0.5), Inches(0.55),
        size=24, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

    # ── RIGHT: Structured output ─────────────────────────────────────────────
    card(slide, Inches(9.38), Inches(1.42), Inches(3.62), Inches(5.65), CARD, GREEN)
    txt(slide, "📤  STRUCTURED OUTPUT", Inches(9.52), Inches(1.55), Inches(3.3), Inches(0.4),
        size=11, bold=True, color=GREEN)

    items_out = [
        ("🎫 ticket", "Fixed login bug INC-231",
         "2h · done · confidence 0.95"),
        ("📅 meeting", "Sprint planning",
         "1h · done · confidence 0.88"),
        ("👁 review", "Data pipeline PR review",
         "— hrs · in_progress · needs clarification"),
        ("🚫 skipped", "Power cut noise",
         "→ filtered, not created"),
    ]
    for j, (cat, task, meta) in enumerate(items_out):
        ty = Inches(2.05 + j * 1.2)
        col_ = GREEN if j < 3 else RED
        rect(slide, Inches(9.45), ty, Inches(3.38), Inches(1.08), CARD2)
        rect(slide, Inches(9.45), ty, Inches(0.04), Inches(1.08), col_)
        txt(slide, cat,  Inches(9.56), ty + Inches(0.05), Inches(3.1), Inches(0.3),
            size=9.5, bold=True, color=col_)
        txt(slide, task, Inches(9.56), ty + Inches(0.33), Inches(3.1), Inches(0.38),
            size=10.5, color=WHITE)
        txt(slide, meta, Inches(9.56), ty + Inches(0.69), Inches(3.1), Inches(0.32),
            size=8.5, color=LGRAY)

    bottom_bar(slide, CARD, "3 work items created · 1 noise item filtered · 1 clarification prompt triggered")


# ═════════════════════════════════════════════════════════════════════════════
#  SLIDE 4 — PERSONAL DASHBOARD (employee features)
# ═════════════════════════════════════════════════════════════════════════════
def s4_dashboard(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, prs)
    top_stripe(slide, GREEN)
    slide_title(slide, "Personal Dashboard", "Every employee gets a full view of their own work")

    badge(slide, "EMPLOYEE VIEW", Inches(10.0), Inches(0.2), GREEN, w_in=2.2)

    # ── KPI row ──────────────────────────────────────────────────────────────
    kpis = [
        ("⏱", "48.5h", "This Week",     BLUE),
        ("✅", "23",    "Tasks Done",    GREEN),
        ("🔄", "8",     "In Progress",   AMBER),
        ("🚧", "3",     "Blocked",       RED),
        ("📋", "34",    "Total Items",   VIOLET),
    ]
    for i, (ico, val, lbl, col) in enumerate(kpis):
        lx = Inches(0.32 + i * 2.59)
        card(slide, lx, Inches(1.42), Inches(2.42), Inches(1.1), CARD, col)
        txt(slide, ico, lx + Inches(0.12), Inches(1.49), Inches(0.5), Inches(0.5),
            size=20, align=PP_ALIGN.CENTER)
        txt(slide, val, lx + Inches(0.65), Inches(1.46), Inches(1.6), Inches(0.55),
            size=24, bold=True, color=col)
        txt(slide, lbl, lx + Inches(0.65), Inches(2.0), Inches(1.6), Inches(0.36),
            size=10, color=LGRAY)

    # ── Charts row (visual mockups) ───────────────────────────────────────────
    # Chart 1: Hours by Category (bar chart mockup)
    card(slide, Inches(0.32), Inches(2.72), Inches(4.3), Inches(3.0), CARD, BLUE)
    txt(slide, "📊  Hours by Category", Inches(0.48), Inches(2.85), Inches(4.0), Inches(0.35),
        size=11, bold=True, color=BLUE)

    cats = [("ticket", 18, BLUE), ("meeting", 9, VIOLET), ("review", 7, GREEN),
            ("admin", 5, AMBER), ("learning", 4, TEAL)]
    max_h = 18
    for j, (lbl, val, col) in enumerate(cats):
        bar_w = Inches(0.5)
        bar_h = Inches((val / max_h) * 1.6)
        bx = Inches(0.6 + j * 0.82)
        by = Inches(5.3) - bar_h
        rect(slide, bx, by, bar_w, bar_h, col)
        txt(slide, str(val), bx, by - Inches(0.28), bar_w, Inches(0.28),
            size=9, bold=True, color=col, align=PP_ALIGN.CENTER)
        txt(slide, lbl, bx - Inches(0.08), Inches(5.32), Inches(0.68), Inches(0.28),
            size=7.5, color=LGRAY, align=PP_ALIGN.CENTER)

    # Chart 2: Status donut (circle mockup)
    card(slide, Inches(4.78), Inches(2.72), Inches(3.0), Inches(3.0), CARD, GREEN)
    txt(slide, "🍩  Status Split", Inches(4.94), Inches(2.85), Inches(2.7), Inches(0.35),
        size=11, bold=True, color=GREEN)
    # Layered ovals for donut effect
    cx, cy = Inches(6.28), Inches(4.35)
    r_out = Inches(0.85)
    r_in  = Inches(0.48)
    status_segs = [(GREEN, 0.55), (AMBER, 0.25), (RED, 0.12), (VIOLET, 0.08)]
    oval(slide, cx - r_out, cy - r_out, r_out*2, r_out*2, GREEN)
    oval(slide, cx - r_out*0.8, cy - r_out*0.8, r_out*1.6, r_out*1.6, AMBER)
    oval(slide, cx - r_out*0.55, cy - r_out*0.55, r_out*1.1, r_out*1.1, RED)
    oval(slide, cx - r_in, cy - r_in, r_in*2, r_in*2, CARD)
    txt(slide, "67%\ndone", cx - r_in, cy - r_in, r_in*2, r_in*2,
        size=9, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    # Legend
    legend = [("Done 67%", GREEN), ("In Progress 25%", AMBER), ("Blocked 8%", RED)]
    for k, (lbl, col) in enumerate(legend):
        ly = Inches(3.15 + k * 0.38)
        rect(slide, Inches(4.94), ly, Inches(0.22), Inches(0.22), col)
        txt(slide, lbl, Inches(5.22), ly, Inches(1.5), Inches(0.25), size=9, color=LGRAY)

    # Chart 3: Daily trend (line mockup)
    card(slide, Inches(7.94), Inches(2.72), Inches(5.07), Inches(3.0), CARD, AMBER)
    txt(slide, "📈  Daily Hours Trend", Inches(8.1), Inches(2.85), Inches(4.7), Inches(0.35),
        size=11, bold=True, color=AMBER)
    # Simple trend bars
    day_vals = [6, 8.5, 7, 9.5, 10, 6.5, 8]
    days_lbl = ["Mon","Tue","Wed","Thu","Fri","Mon","Tue"]
    max_v = 10
    for j, (val, dlbl) in enumerate(zip(day_vals, days_lbl)):
        bx = Inches(8.15 + j * 0.68)
        bar_h = Inches((val / max_v) * 1.55)
        by = Inches(5.3) - bar_h
        col_b = AMBER if val < 9 else GREEN
        rect(slide, bx, by, Inches(0.5), bar_h, col_b)
        txt(slide, str(val), bx, by - Inches(0.28), Inches(0.5), Inches(0.28),
            size=8, bold=True, color=col_b, align=PP_ALIGN.CENTER)
        txt(slide, dlbl, bx, Inches(5.32), Inches(0.5), Inches(0.28),
            size=7.5, color=LGRAY, align=PP_ALIGN.CENTER)

    # Feature badges bottom
    feats = ["Date range filter", "Inline editing", "Sortable table", "Category filter"]
    for i, f in enumerate(feats):
        rect(slide, Inches(0.32 + i * 3.27), Inches(5.88), Inches(3.05), Inches(0.3), CARD2)
        txt(slide, f"✓  {f}", Inches(0.42 + i * 3.27), Inches(5.88), Inches(2.85), Inches(0.3),
            size=9.5, color=GREEN)

    bottom_bar(slide, CARD, "All charts are interactive (Plotly) · Date range: this week / last 7 / last 30 / custom")


# ═════════════════════════════════════════════════════════════════════════════
#  SLIDE 5 — TEAM DASHBOARD (manager view)
# ═════════════════════════════════════════════════════════════════════════════
def s5_team(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, prs)
    top_stripe(slide, VIOLET)
    slide_title(slide, "Team Dashboard", "Managers get a real-time view across the entire team")

    badge(slide, "MANAGER / ADMIN", Inches(9.5), Inches(0.2), VIOLET, w_in=2.7)

    # ── Stacked bar mockup (main chart) ──────────────────────────────────────
    card(slide, Inches(0.32), Inches(1.42), Inches(6.2), Inches(3.6), CARD, VIOLET)
    txt(slide, "📊  Team Hours — Stacked by Employee",
        Inches(0.48), Inches(1.56), Inches(5.8), Inches(0.38),
        size=11, bold=True, color=VIOLET)

    emp_colors = [BLUE, GREEN, AMBER, TEAL, VIOLET]
    emp_names  = ["Alice", "Bob", "Carol", "Dave", "Eve"]
    day_data   = [
        [6,5,4,7,3],[7,4,6,5,5],[5,6,7,4,6],[8,5,5,6,4],[7,6,4,7,5]
    ]
    days_x = ["Mon","Tue","Wed","Thu","Fri"]
    bar_w = Inches(0.52)

    for j, (day, day_vals) in enumerate(zip(days_x, day_data)):
        base = Inches(4.75)
        bx = Inches(1.0 + j * 1.02)
        for k, (val, col) in enumerate(zip(day_vals, emp_colors)):
            bh = Inches(val * 0.16)
            base -= bh
            rect(slide, bx, base, bar_w, bh, col)
        txt(slide, day, bx, Inches(4.8), bar_w, Inches(0.28),
            size=8, color=LGRAY, align=PP_ALIGN.CENTER)

    # Legend
    for k, (name, col) in enumerate(zip(emp_names, emp_colors)):
        lx = Inches(1.0 + k * 1.05)
        rect(slide, lx, Inches(2.0), Inches(0.18), Inches(0.18), col)
        txt(slide, name, lx + Inches(0.24), Inches(1.97), Inches(0.75), Inches(0.25),
            size=8.5, color=LGRAY)

    # ── Right: Blocked items panel ────────────────────────────────────────────
    card(slide, Inches(6.7), Inches(1.42), Inches(6.31), Inches(3.6), CARD, RED)
    txt(slide, "🚧  Blocked Items — Needs Attention",
        Inches(6.86), Inches(1.56), Inches(5.9), Inches(0.38),
        size=11, bold=True, color=RED)

    blocked = [
        ("Dave M.", "Waiting for DB access from IT", "ticket", "2 days"),
        ("Carol S.", "API credentials not received",  "project","1 day"),
        ("Bob K.",  "Blocked on legal review",        "admin",  "3 days"),
    ]
    for j, (emp, issue, cat, age) in enumerate(blocked):
        ty = Inches(2.05 + j * 0.95)
        rect(slide, Inches(6.78), ty, Inches(6.05), Inches(0.85), CARD2)
        rect(slide, Inches(6.78), ty, Inches(0.04), Inches(0.85), RED)
        txt(slide, emp,   Inches(6.9),  ty + Inches(0.06), Inches(1.5), Inches(0.3),
            size=11, bold=True, color=WHITE)
        txt(slide, issue, Inches(6.9),  ty + Inches(0.38), Inches(4.0), Inches(0.4),
            size=10, color=LGRAY)
        badge(slide, age, Inches(12.25), ty + Inches(0.26), RED, w_in=0.7, h_in=0.27, size=9)

    # ── Per-employee summary cards ────────────────────────────────────────────
    card(slide, Inches(0.32), Inches(5.2), Inches(12.69), Inches(1.6), CARD, VIOLET)
    txt(slide, "👥  Weekly Summary — Per Employee",
        Inches(0.48), Inches(5.32), Inches(5.0), Inches(0.35),
        size=11, bold=True, color=VIOLET)

    emp_stats = [
        ("Alice", 42.0, 12, 2, 0),
        ("Bob",   38.5, 10, 3, 1),
        ("Carol", 45.0, 14, 1, 2),
        ("Dave",  35.0, 9,  4, 1),
    ]
    for i, (name, hrs, done, ip, blk) in enumerate(emp_stats):
        lx = Inches(0.42 + i * 3.18)
        rect(slide, lx, Inches(5.72), Inches(2.95), Inches(0.88), CARD2)
        txt(slide, name, lx + Inches(0.12), Inches(5.78), Inches(1.3), Inches(0.32),
            size=11, bold=True, color=WHITE)
        txt(slide, f"{hrs}h", lx + Inches(1.5), Inches(5.78), Inches(1.1), Inches(0.32),
            size=13, bold=True, color=VIOLET, align=PP_ALIGN.RIGHT)
        txt(slide, f"✅{done}  🔄{ip}  🚧{blk}",
            lx + Inches(0.12), Inches(6.12), Inches(2.7), Inches(0.3),
            size=9.5, color=LGRAY)

    bottom_bar(slide, CARD,
               "Filter by employee · category · date range · status  ·  "
               "All aggregations enforced server-side per role")


# ═════════════════════════════════════════════════════════════════════════════
#  SLIDE 6 — CHAT ASSISTANT (LangGraph tool flow)
# ═════════════════════════════════════════════════════════════════════════════
def s6_chat(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, prs)
    top_stripe(slide, TEAL)
    slide_title(slide, "AI Chat Assistant",
                "LangGraph ReAct agent · 3 tools · hybrid SQL + semantic retrieval")

    # ── User question box ─────────────────────────────────────────────────────
    card(slide, Inches(0.32), Inches(1.42), Inches(3.6), Inches(1.1), CARD, TEAL)
    txt(slide, "💬  USER ASKS", Inches(0.48), Inches(1.5), Inches(3.3), Inches(0.3),
        size=10, bold=True, color=TEAL)
    txt(slide, '"How many hours did I log on tickets last week?"',
        Inches(0.42), Inches(1.82), Inches(3.4), Inches(0.62),
        size=11, color=WHITE, italic=True)

    # Arrow down
    arrow_down(slide, Inches(1.9), Inches(2.58))

    # ── Agent box ────────────────────────────────────────────────────────────
    card(slide, Inches(0.32), Inches(3.0), Inches(3.6), Inches(0.9), CARD, VIOLET)
    txt(slide, "🧠  LangGraph ReAct Agent",
        Inches(0.48), Inches(3.08), Inches(3.3), Inches(0.35),
        size=11, bold=True, color=VIOLET)
    txt(slide, "Claude Sonnet 4.6 · temp=0 · session history (last 10 turns)",
        Inches(0.48), Inches(3.45), Inches(3.3), Inches(0.38),
        size=9.5, color=LGRAY)

    # Arrows to tools
    for xi in [Inches(2.5), Inches(6.2), Inches(9.9)]:
        txt(slide, "➜", Inches(3.98), Inches(3.3), Inches(0.4), Inches(0.4),
            size=16, bold=True, color=VIOLET, align=PP_ALIGN.CENTER)

    txt(slide, "➜", Inches(3.95), Inches(3.3), Inches(0.5), Inches(0.4),
        size=18, bold=True, color=VIOLET, align=PP_ALIGN.CENTER)

    # ── Three tool cards ─────────────────────────────────────────────────────
    tools = [
        ("📅", "date_resolver", TEAL,
         "Pure Python\nZero LLM calls",
         ['"last week"', '"yesterday"', '"this month"', '"last 30 days"'],
         "→ {start_date, end_date}"),
        ("🗄", "sql_query", BLUE,
         "SQLite via\ndashboard_service",
         ["total_hours_summary", "hours_by_category", "status_distribution",
          "daily_trend", "list_items", "team_summary (manager)"],
         "→ Counts · totals · lists"),
        ("🔍", "vector_search", VIOLET,
         "ChromaDB semantic\ntext-embedding-3-small",
         ['"what was I working on…"', '"find tasks about API"',
          "Date filter via work_date_num", "Up to 20 results returned"],
         "→ Fuzzy semantic recall"),
    ]

    for i, (ico, name, col, sub, features, ret) in enumerate(tools):
        lx = Inches(4.55 + i * 2.95)
        card(slide, lx, Inches(1.42), Inches(2.8), Inches(5.65), CARD, col)

        # Icon + name
        txt(slide, ico, lx + Inches(0.15), Inches(1.5), Inches(0.5), Inches(0.5),
            size=20, align=PP_ALIGN.CENTER)
        txt(slide, name, lx + Inches(0.7), Inches(1.52), Inches(1.9), Inches(0.4),
            size=11, bold=True, color=col)
        txt(slide, sub, lx + Inches(0.15), Inches(1.95), Inches(2.55), Inches(0.5),
            size=9, color=LGRAY, italic=True)

        hbar(slide, Inches(2.5), color=DGRAY)

        for j, feat in enumerate(features):
            txt(slide, f"• {feat}", lx + Inches(0.15), Inches(2.6 + j * 0.55),
                Inches(2.55), Inches(0.5), size=9.5, color=LGRAY)

        # Return box
        rect(slide, lx + Inches(0.15), Inches(6.4), Inches(2.55), Inches(0.3), col)
        txt(slide, ret, lx + Inches(0.15), Inches(6.4), Inches(2.55), Inches(0.3),
            size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # ── Query type + example answer ───────────────────────────────────────────
    card(slide, Inches(0.32), Inches(4.05), Inches(3.6), Inches(1.55), CARD, GREEN)
    txt(slide, "⚡  HYBRID ANSWER", Inches(0.48), Inches(4.12), Inches(3.3), Inches(0.3),
        size=10, bold=True, color=GREEN)
    txt(slide, '"You logged 18.5h on ticket tasks last week (Mon–Fri). '
               'Top tickets: INC-231 (4h), PROJ-88 (3.5h)…"',
        Inches(0.42), Inches(4.45), Inches(3.4), Inches(0.9),
        size=10, color=WHITE, italic=True)

    # Query type indicators
    types = [("🗄 SQL", BLUE, "counts/totals"), ("🔍 Vector", VIOLET, "semantic"), ("⚡ Hybrid", GREEN, "both")]
    for i, (label, col, note) in enumerate(types):
        lx = Inches(0.38 + i * 1.22)
        rect(slide, lx, Inches(5.72), Inches(1.1), Inches(0.55), col)
        txt(slide, label, lx, Inches(5.72), Inches(1.1), Inches(0.3),
            size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        txt(slide, note, lx, Inches(6.0), Inches(1.1), Inches(0.25),
            size=8, color=LGRAY, align=PP_ALIGN.CENTER)

    bottom_bar(slide, CARD,
               "Source citations shown per answer · Session history (10 turns) · "
               "Role-aware: employee sees own data · manager sees team")


# ═════════════════════════════════════════════════════════════════════════════
#  SLIDE 7 — ROLE-BASED ACCESS (visual 3 columns)
# ═════════════════════════════════════════════════════════════════════════════
def s7_rbac(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, prs)
    top_stripe(slide, AMBER)
    slide_title(slide, "Role-Based Access Control",
                "Server-enforced at every route · service · query · JWT carries user_id + role")

    roles = [
        ("👤", "EMPLOYEE", BLUE,
         [("📝", "Submit & re-submit daily updates"),
          ("🤖", "View AI preview, edit before saving"),
          ("🗑", "Soft-delete own work logs"),
          ("✏️", "Inline-edit own items only"),
          ("📊", "Personal Dashboard — KPIs + charts"),
          ("💬", "Chat: own data only"),
          ("📜", "Own chat session history"),
          ("⚠️", "Cannot see other employees' data"),
          ]),
        ("👥", "MANAGER", VIOLET,
         [("✅", "All Employee capabilities"),
          ("📊", "Team Dashboard — stacked bars"),
          ("🚧", "Blocked items panel (team-wide)"),
          ("🗂", "Filter team by employee/category/date"),
          ("💬", "Chat: query across entire team"),
          ("🏢", "Org hierarchy via manager_id"),
          ("📈", "Team aggregations & trends"),
          ("🚫", "Cannot edit other employees' items"),
          ]),
        ("🔑", "ADMIN", GREEN,
         [("✅", "All Manager capabilities"),
          ("👁", "View all users (any team/dept)"),
          ("🛠", "Update role · team · department"),
          ("🔓", "Activate / deactivate accounts"),
          ("🛡", "Cannot change own role or deactivate self"),
          ("🔬", "Extraction error queue"),
          ("♻️", "ChromaDB reindex from SQLite"),
          ("🌱", "Seed 16 users + 30 days of data"),
          ]),
    ]

    for i, (ico, role_name, col, items) in enumerate(roles):
        lx = Inches(0.32 + i * 4.35)
        card(slide, lx, Inches(1.42), Inches(4.15), Inches(5.68), CARD, col)

        # Large role icon
        oval(slide, lx + Inches(1.55), Inches(1.58), Inches(1.05), Inches(1.05), col)
        txt(slide, ico, lx + Inches(1.55), Inches(1.62), Inches(1.05), Inches(1.0),
            size=26, align=PP_ALIGN.CENTER)

        txt(slide, role_name, lx + Inches(0.15), Inches(2.72), Inches(3.85), Inches(0.42),
            size=16, bold=True, color=col, align=PP_ALIGN.CENTER)

        hbar(slide, Inches(3.2), color=DGRAY)

        for j, (item_ico, item_txt) in enumerate(items):
            ty = Inches(3.3 + j * 0.6)
            txt(slide, item_ico, lx + Inches(0.18), ty, Inches(0.38), Inches(0.5),
                size=13, align=PP_ALIGN.CENTER)
            color_ = LGRAY if "Cannot" not in item_txt and "All" not in item_txt else (
                     RED if "Cannot" in item_txt else GREEN)
            txt(slide, item_txt, lx + Inches(0.62), ty + Inches(0.06), Inches(3.38), Inches(0.46),
                size=10, color=color_)

    bottom_bar(slide, CARD,
               "Inactive users blocked at JWT validation  ·  "
               "user_id set server-side from JWT — never from LLM output  ·  "
               "Role guards enforced in every service layer")


# ═════════════════════════════════════════════════════════════════════════════
#  SLIDE 8 — ARCHITECTURE & TECH STACK
# ═════════════════════════════════════════════════════════════════════════════
def s8_arch(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, prs)
    top_stripe(slide, TEAL)
    slide_title(slide, "Architecture & Tech Stack",
                "FastAPI · LangChain · LangGraph · Streamlit · SQLite · ChromaDB")

    # ── Left: Layered architecture diagram ────────────────────────────────────
    layers = [
        ("🖥",  "Streamlit Frontend",       TEAL,
         "Login · Submit · My Dashboard · Team Dashboard · Chat · Admin"),
        ("⚡",  "FastAPI Backend  :8000",   BLUE,
         "6 Routers · JWT + bcrypt auth · Pydantic validation · async"),
        ("🧠",  "LangChain / LangGraph",    VIOLET,
         "Extraction chain  +  ReAct agent (3 tools)  +  OutputFixingParser"),
        ("🗄",  "SQLite  +  ChromaDB",      GREEN,
         "ORM via SQLAlchemy 2.0  ·  text-embedding-3-small  ·  --reindex CLI"),
    ]

    for i, (ico, name, col, desc) in enumerate(layers):
        ty = Inches(1.42 + i * 1.35)
        card(slide, Inches(0.32), ty, Inches(7.8), Inches(1.2), CARD, col)
        txt(slide, ico, Inches(0.42), ty + Inches(0.28), Inches(0.55), Inches(0.6),
            size=22, align=PP_ALIGN.CENTER)
        txt(slide, name, Inches(1.1), ty + Inches(0.1), Inches(6.7), Inches(0.42),
            size=13, bold=True, color=col)
        txt(slide, desc, Inches(1.1), ty + Inches(0.52), Inches(6.7), Inches(0.58),
            size=10.5, color=LGRAY)
        # Connector arrow (except last)
        if i < 3:
            txt(slide, "↕", Inches(3.7), ty + Inches(1.2), Inches(0.4), Inches(0.18),
                size=10, color=DGRAY, align=PP_ALIGN.CENTER)

    # ── Right: Model config cards ─────────────────────────────────────────────
    txt(slide, "LLM Configuration", Inches(8.4), Inches(1.42), Inches(4.6), Inches(0.38),
        size=13, bold=True, color=TEAL)

    llm_rows = [
        ("Extraction",  "Claude Sonnet 4.6", "temp=0 · accuracy", BLUE),
        ("JSON Fix",    "Claude Haiku 4.5",   "temp=0 · speed",    GREEN),
        ("Chat Agent",  "Claude Sonnet 4.6", "temp=0 · reasoning", VIOLET),
        ("Alt models",  "GPT-4o / GPT-5 / o4-mini-gs", "swappable via .env", AMBER),
        ("Embeddings",  "text-embedding-3-small", "OpenAI → ChromaDB", TEAL),
    ]
    for j, (role, model, note, col) in enumerate(llm_rows):
        ty = Inches(1.88 + j * 0.82)
        rect(slide, Inches(8.4), ty, Inches(4.6), Inches(0.72), CARD2)
        rect(slide, Inches(8.4), ty, Inches(0.04), Inches(0.72), col)
        txt(slide, role, Inches(8.52), ty + Inches(0.06), Inches(1.1), Inches(0.3),
            size=9.5, bold=True, color=col)
        txt(slide, model, Inches(9.68), ty + Inches(0.06), Inches(2.0), Inches(0.3),
            size=9.5, color=WHITE)
        txt(slide, note, Inches(9.68), ty + Inches(0.36), Inches(3.0), Inches(0.3),
            size=8.5, color=LGRAY, italic=True)

    # Metric badges
    metrics = [("218", "Tests", BLUE), ("6", "Phases", GREEN),
               ("10", "Categories", VIOLET), ("8h", "JWT expiry", AMBER)]
    for i, (val, lbl, col) in enumerate(metrics):
        lx = Inches(8.4 + i * 1.16)
        ty = Inches(6.27)
        rect(slide, lx, ty, Inches(1.06), Inches(0.82), CARD2)
        rect(slide, lx, ty, Inches(1.06), Inches(0.04), col)
        txt(slide, val, lx, ty + Inches(0.06), Inches(1.06), Inches(0.42),
            size=20, bold=True, color=col, align=PP_ALIGN.CENTER)
        txt(slide, lbl, lx, ty + Inches(0.5), Inches(1.06), Inches(0.26),
            size=8.5, color=LGRAY, align=PP_ALIGN.CENTER)

    bottom_bar(slide, CARD,
               "Azure Bedrock proxy (Claude)  ·  NLP API (GPT)  ·  "
               "Enterprise TLS via cacert.pem  ·  SQLite → Postgres migration planned")


# ═════════════════════════════════════════════════════════════════════════════
#  SLIDE 9 — VALUE & ROADMAP
# ═════════════════════════════════════════════════════════════════════════════
def s9_roadmap(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, prs)
    top_stripe(slide, GREEN)
    slide_title(slide, "Value & Roadmap",
                "Prototype complete · All 6 phases delivered · Production path defined")

    # ── Left: Value propositions (icon cards) ─────────────────────────────────
    txt(slide, "Why It Matters", Inches(0.32), Inches(1.42), Inches(6.2), Inches(0.38),
        size=14, bold=True, color=GREEN)

    values = [
        ("📝", "Zero Friction",      BLUE,   "Plain English input — no forms, no dropdowns"),
        ("🤖", "AI + Human Control", VIOLET, "LLM extracts; human confirms — accuracy without rigidity"),
        ("👁",  "Live Visibility",   GREEN,  "Managers see team blockers & trends in real time"),
        ("💬", "Ask Anything",       TEAL,   '"What was I blocked on last week?" → cited answer'),
        ("📋", "Full Audit Trail",   AMBER,  "Immutable raw_message · correction flags · model metadata"),
        ("🔒", "Zero-trust Data",    RED,    "Server-side identity · role guards · bcrypt + JWT"),
    ]
    for i, (ico, title, col, desc) in enumerate(values):
        row, col_i = divmod(i, 2)
        lx = Inches(0.32 + col_i * 3.12)
        ty = Inches(1.9 + row * 1.15)
        card(slide, lx, ty, Inches(2.95), Inches(1.0), CARD, col)
        txt(slide, ico, lx + Inches(0.12), ty + Inches(0.22), Inches(0.5), Inches(0.55),
            size=20, align=PP_ALIGN.CENTER)
        txt(slide, title, lx + Inches(0.68), ty + Inches(0.08), Inches(2.12), Inches(0.35),
            size=12, bold=True, color=col)
        txt(slide, desc, lx + Inches(0.68), ty + Inches(0.45), Inches(2.12), Inches(0.48),
            size=9.5, color=LGRAY)

    # ── Right: Phase timeline ─────────────────────────────────────────────────
    txt(slide, "Delivery Phases", Inches(6.9), Inches(1.42), Inches(6.1), Inches(0.38),
        size=14, bold=True, color=VIOLET)

    phases = [
        ("✅", "Phase 1–2", "Scaffolding · Auth · LLM Extraction · ChromaDB",     GREEN,  True),
        ("✅", "Phase 3–4", "Dashboards · CRUD · LangGraph Chat Agent",             GREEN,  True),
        ("✅", "Phase 5–6", "Team Dashboard · Admin · Seed Data · Polish",          GREEN,  True),
        ("▶",  "Next",      "PostgreSQL migration · Slack/Teams bot integration",   BLUE,   False),
        ("▶",  "Future",    "OAuth/SSO · Manager approvals · Jira sync · Digests", VIOLET, False),
    ]

    # Timeline line
    rect(slide, Inches(7.55), Inches(1.9), Inches(0.04), Inches(5.0), DGRAY)

    for i, (ico, phase, desc, col, done) in enumerate(phases):
        ty = Inches(1.88 + i * 0.98)
        # Dot on timeline
        oval(slide, Inches(7.38), ty + Inches(0.12), Inches(0.38), Inches(0.38), col)
        txt(slide, ico, Inches(7.38), ty + Inches(0.1), Inches(0.38), Inches(0.38),
            size=12, align=PP_ALIGN.CENTER, color=WHITE)
        # Content
        card(slide, Inches(7.92), ty, Inches(5.06), Inches(0.84), CARD, col if not done else GREEN)
        txt(slide, phase, Inches(8.06), ty + Inches(0.06), Inches(1.2), Inches(0.34),
            size=11, bold=True, color=col if not done else GREEN)
        txt(slide, desc, Inches(9.3), ty + Inches(0.06), Inches(3.55), Inches(0.72),
            size=10, color=LGRAY if done else WHITE)

    # CTA bottom
    rect(slide, 0, H - Inches(0.7), W, Inches(0.7), VIOLET)
    txt(slide, "WorkTrack AI  —  Prototype Complete  ·  Ready for Review",
        Inches(0.5), H - Inches(0.65), Inches(12.3), Inches(0.58),
        size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ═════════════════════════════════════════════════════════════════════════════
#  BUILD
# ═════════════════════════════════════════════════════════════════════════════
def build():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    s1_cover(prs)
    s2_pipeline(prs)
    s3_extraction(prs)
    s4_dashboard(prs)
    s5_team(prs)
    s6_chat(prs)
    s7_rbac(prs)
    s8_arch(prs)
    s9_roadmap(prs)

    out = "/home/x288712/WorkTrack AI/WorkTrack_AI_Presentation.pptx"
    prs.save(out)
    print(f"✓ Saved → {out}  ({len(prs.slides)} slides)")

if __name__ == "__main__":
    build()
